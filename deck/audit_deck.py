#!/usr/bin/env python3
"""
audit_deck.py — geometry and content QA for the built .pptx.

    python deck/audit_deck.py [deck.pptx]

LibreOffice cannot render in every environment, so this reads the real shape
geometry out of the file instead and reports the defects that would otherwise
only surface visually.

It deliberately reasons about where text actually *lands*, not where its box is.
A right-aligned label in a wide box occupies only its right end, so comparing raw
box rectangles produces a flood of overlaps that do not exist on screen — the
footer pair on every slide being the obvious case. Estimated text extent is used
for the margin and overlap checks; the full shape box is still used for the
off-slide check, since a box hanging past the edge is worth knowing about even
when its text is short.

Reported:
  * shapes whose box leaves the slide
  * text whose rendered extent leaves the slide or crowds the edge
  * text that will not fit its box at the declared size
  * text that visually collides with other text
  * placeholder text
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

DECK = Path(sys.argv[1]) if len(sys.argv) > 1 else \
    Path(__file__).with_name("Kingdom_of_Ultra_Secure_Collaboration_Storyboard.pptx")

EDGE_MIN = 0.12        # text should not land closer than this to a slide edge
CHAR_W = 0.50          # mean glyph advance as a fraction of font size
LINE_H = 1.22
PICTURE = 13


def inches(v):
    return None if v is None else Emu(v).inches


def paras(shape):
    out = []
    for p in shape.text_frame.paragraphs:
        size = next((r.font.size.pt for r in reversed(p.runs) if r.font.size), None)
        align = p.alignment
        text = "".join(r.text for r in p.runs)
        if text.strip():
            out.append((text, size, align))
    return out


def text_extent(shape, x, y, w, h):
    """Estimated rectangle the glyphs actually occupy, honouring alignment."""
    ps = paras(shape)
    if not ps:
        return None
    widest, total_h = 0.0, 0.0
    for text, pt, _ in ps:
        pt = pt or 14
        cw = (pt * CHAR_W) / 72.0
        line_w = len(text) * cw
        per_line = max(1, int(w / cw)) if cw else 1
        lines = max(1, -(-len(text) // per_line))
        widest = max(widest, min(line_w, w))
        total_h += lines * (pt * LINE_H) / 72.0
    align = ps[0][2]
    if align == PP_ALIGN.RIGHT:
        tx = x + w - widest
    elif align == PP_ALIGN.CENTER:
        tx = x + (w - widest) / 2
    else:
        tx = x
    return tx, y, widest, min(total_h, h), total_h


def main():
    if not DECK.exists():
        sys.exit(f"no such deck: {DECK}")
    prs = Presentation(str(DECK))
    W, H = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    print(f"  {DECK.name}")
    print(f"  {len(prs.slides)} slides, {W:.2f}in x {H:.2f}in\n")

    hard, soft = [], []

    for i, slide in enumerate(prs.slides, 1):
        spans = []
        for shape in slide.shapes:
            x, y, w, h = (inches(shape.left), inches(shape.top),
                          inches(shape.width), inches(shape.height))
            if None in (x, y, w, h):
                continue

            # box past the slide edge — cosmetic when the text is short, but real
            if x + w > W + 0.01 or y + h > H + 0.01 or x < -0.01 or y < -0.01:
                over = max(x + w - W, y + h - H, -x, -y)
                soft.append(f"slide {i}: {shape.shape_type} box extends {over:.2f}\" past the slide "
                            f"(at {x:.2f},{y:.2f} size {w:.2f}x{h:.2f})")

            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            label = shape.text_frame.text.strip().replace("\n", " ")[:30]

            for bad in ("lorem", "ipsum", "todo", "[insert"):
                if bad in shape.text_frame.text.lower():
                    hard.append(f"slide {i}: placeholder text '{bad}' — '{label}'")

            ext = text_extent(shape, x, y, w, h)
            if not ext:
                continue
            tx, ty, tw, th, need = ext

            if need > h * 1.4:
                hard.append(f"slide {i}: text overflows its box "
                            f"(needs ~{need:.2f}\" of {h:.2f}\") — '{label}'")
            if tx < EDGE_MIN or ty < EDGE_MIN or tx + tw > W - EDGE_MIN or ty + th > H - EDGE_MIN:
                hard.append(f"slide {i}: text lands within {EDGE_MIN}\" of a slide edge — '{label}'")
            spans.append((tx, ty, tw, th, label))

        for a in range(len(spans)):
            for b in range(a + 1, len(spans)):
                ax, ay, aw, ah, at = spans[a]
                bx, by, bw, bh, bt = spans[b]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.08 and oy > 0.08:
                    hard.append(f"slide {i}: text collides — '{at}' / '{bt}' "
                                f"({ox:.2f}x{oy:.2f}in)")

    print(f"  {len(hard)} user-visible issue(s)")
    for s in hard:
        print("   -", s)
    print(f"\n  {len(soft)} cosmetic note(s) (box past the edge; short text may still look fine)")
    for s in soft[:6]:
        print("   -", s)
    if len(soft) > 6:
        print(f"   … and {len(soft)-6} more of the same kind")
    return 1 if hard else 0


if __name__ == "__main__":
    sys.exit(main())
