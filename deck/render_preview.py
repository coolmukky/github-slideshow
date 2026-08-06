#!/usr/bin/env python3
"""
render_preview.py — render slides of a .pptx to HTML for visual review.

    python deck/render_preview.py [deck.pptx] [out.html] [slide numbers...]

LibreOffice's Impress component is not installed in every environment (it fails
here even on a plain text file), so the usual convert-to-PDF visual pass is
unavailable. This reads the real shape geometry, text and embedded images out of
the .pptx with python-pptx and lays them out at true scale in HTML, which a
browser can then screenshot.

It renders what is actually in the file — pictures, autoshape fills, and every
text run's size, weight, colour and alignment — so it catches generator bugs
rather than reflecting assumptions about them. It is an approximation of
PowerPoint's own layout, not a substitute for it: line breaking and font metrics
will differ slightly. Treat it as a layout check, not a pixel proof.
"""
import base64
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

PPI = 96  # css px per inch

DECK = Path(sys.argv[1]) if len(sys.argv) > 1 else \
    Path(__file__).with_name("Kingdom_of_Ultra_Secure_Collaboration_Storyboard.pptx")
OUT = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("/tmp/deck_preview.html")
WANT = [int(a) for a in sys.argv[3:]] if len(sys.argv) > 3 else None

ALIGN = {PP_ALIGN.LEFT: "left", PP_ALIGN.RIGHT: "right",
         PP_ALIGN.CENTER: "center", PP_ALIGN.JUSTIFY: "justify"}


def px(v):
    return Emu(v).inches * PPI


def rgb(c):
    try:
        if c and c.type is not None and c.rgb is not None:
            return "#%s" % str(c.rgb)
    except Exception:
        pass
    return None


def esc(s):
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def shape_html(sh):
    try:
        x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
    except Exception:
        return ""
    base = f"position:absolute;left:{x:.1f}px;top:{y:.1f}px;width:{w:.1f}px;height:{h:.1f}px;"
    out = []

    # picture
    if sh.shape_type == 13 and getattr(sh, "image", None):
        b64 = base64.b64encode(sh.image.blob).decode()
        out.append(f'<img style="{base}object-fit:cover" '
                   f'src="data:{sh.image.content_type};base64,{b64}">')
        return "".join(out)

    # autoshape fill. Alpha matters: the deck lays translucent scrims over full-bleed
    # images, and drawing those opaque hides the artwork entirely.
    fill_css = ""
    try:
        if sh.fill.type is not None and sh.fill.type == 1:  # solid
            c = rgb(sh.fill.fore_color)
            if c:
                alpha = 1.0
                el = sh.fill.fore_color._xFill.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
                if el is not None:
                    a = el.find("{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
                    if a is not None and a.get("val"):
                        alpha = int(a.get("val")) / 100000.0
                r, g, b = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
                fill_css = f"background:rgba({r},{g},{b},{alpha:.3f});"
    except Exception:
        pass
    line_css = ""
    try:
        c = rgb(sh.line.color)
        if c and sh.line.width:
            line_css = f"border:{max(1, px(sh.line.width)):.0f}px solid {c};"
    except Exception:
        pass
    radius = "border-radius:8px;" if "ROUNDED" in str(sh.shape_type) else ""
    if fill_css or line_css:
        out.append(f'<div style="{base}{fill_css}{line_css}{radius}box-sizing:border-box"></div>')

    # text
    if sh.has_text_frame and sh.text_frame.text.strip():
        paras = []
        for p in sh.text_frame.paragraphs:
            if not p.runs:
                continue
            al = ALIGN.get(p.alignment, "left")
            runs = []
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else 14
                col = rgb(r.font.color) or "#FFFFFF"
                fw = "700" if r.font.bold else "400"
                fam = r.font.name or "Inter"
                sp = f"letter-spacing:{r.font._rPr.get('spc', 0)}" if False else ""
                runs.append(f'<span style="font-size:{sz}px;color:{col};font-weight:{fw};'
                            f'font-family:\'{fam}\',Arial,sans-serif;{sp}">{esc(r.text)}</span>')
            paras.append(f'<p style="margin:0;text-align:{al};line-height:1.18">{"".join(runs)}</p>')
        if paras:
            out.append(f'<div style="{base}padding:2px 4px;box-sizing:border-box;'
                       f'overflow:visible">{"".join(paras)}</div>')
    return "".join(out)


def main():
    prs = Presentation(str(DECK))
    W, H = px(prs.slide_width), px(prs.slide_height)
    html = ['<!doctype html><meta charset=utf-8>',
            '<body style="margin:0;background:#11161f;font-family:Arial,sans-serif">']
    for i, slide in enumerate(prs.slides, 1):
        if WANT and i not in WANT:
            continue
        bg = "#07182D"
        try:
            c = rgb(slide.background.fill.fore_color)
            if c:
                bg = c
        except Exception:
            pass
        html.append(f'<div style="color:#8fa6c4;font:12px monospace;padding:8px 4px 3px">slide {i}</div>')
        html.append(f'<div style="position:relative;width:{W:.0f}px;height:{H:.0f}px;'
                    f'background:{bg};overflow:hidden;outline:1px solid #2a3b57">')
        for sh in slide.shapes:
            html.append(shape_html(sh))
        html.append("</div>")
    html.append("</body>")
    OUT.write_text("".join(html), encoding="utf-8")
    n = len(WANT) if WANT else len(prs.slides)
    print(f"  rendered {n} slide(s) -> {OUT}  ({W:.0f}x{H:.0f}px each)")


if __name__ == "__main__":
    main()
